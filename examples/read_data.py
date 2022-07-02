"""
make slides to present DeepEyeNet examples
https://www.geeksforgeeks.org/creating-and-updating-powerpoint-presentations-in-python-using-python-pptx/
https://python-pptx.readthedocs.io/en/latest/user/placeholders-using.html
"""
from ast import keyword
import json
from pptx import Presentation
from pptx.util import Inches 


def main():
    """
    create an example slides
    """
    # train: 9428, val: 3142, test: 3140
    datapath_list = ['DeepEyeNet_train.json', 'DeepEyeNet_valid.json', 'DeepEyeNet_test.json']
    labels = []
    for datapath in datapath_list:
        with open(datapath) as f:
            data = json.load(f)
            # print(data)
            # Giving Image path 
            
            # Creating an Presentation object
            # ppt = Presentation() 
            
            # # Selecting blank slide
            # blank_slide_layout = ppt.slide_layouts[8] 
            
            # For margins
            # left = top = Inches(1)

            for i, example in enumerate(data):
                for k, v in example.items():
                    print(k, v)
            # if i > 300:
            #     break
            # for k, v in example.items():  # k: img path, v: caption
            #     # Attaching slide to ppt
            #     slide = ppt.slides.add_slide(blank_slide_layout) 
            #     title_placeholder = slide.shapes.title
            #     title_placeholder.text = k  # img path
            #     # adding images
            #     img_placeholder = slide.placeholders[1]
            #     pic = img_placeholder.insert_picture(k)

            #     slide.placeholders[2].text = f'{v}'
                
        # save file
        # ppt.save('deepeyenet_examples.pptx')
        
        # print("Done")


def build_label_dict():
    """
    build label's label2idx and idx2label dicts
    build 3 splits sets of keywords, then intersection, no, separated keyword can't work
    use the whole keywords as one label
    7/1: consider filter by word frequency of > 8
    """
    datapath_list = ['DeepEyeNet_train.json', 'DeepEyeNet_valid.json', 'DeepEyeNet_test.json']
    all_labels = []
    label2idx = {}
    idx2label = {}
    for datapath in datapath_list:
        split_labels = []
        with open(datapath) as f:
            data = json.load(f)
            for i, example in enumerate(data):
                for k, v in example.items():
                    keywords = v['keywords']
                    # print(keywords)
                    split_labels.append(keywords)
        all_labels.append(set(split_labels))

    inter = set.intersection(*all_labels)
    for i, label in enumerate(inter):
        label2idx[label] = i
        idx2label[i] = label
    
    with open(f'intersection_label2idx.json', 'w') as f1:
        json.dump(label2idx, f1)
    
    with open(f'intersection_idx2label.json', 'w') as f2:
        json.dump(idx2label, f2)


def filter_zero_classes():
    """
    filter out train/val/test's zero count classes
    1. get 3 splits' zero classes and build a new nonzero classes dict
    2. use nonzero class dict to filter 3 splits and save 3 splits to new json
    """
    with open(f'intersection_label2idx.json', 'r') as f1:
        label2idx = json.load(f1)  # full dict
    with open(f'intersection_idx2label.json', 'r') as f2:
        idx2label = json.load(f2)  # full dict

    datapath_list = ['DeepEyeNet_train.json', 'DeepEyeNet_valid.json', 'DeepEyeNet_test.json']
    all_indices = list(label2idx.values())
    all_split_zero_classes = []
    for i, file in enumerate(datapath_list):
        split_nonzero_classes = []
        with open(file) as f:
            split_list = json.load(f)
            for example in split_list:
                for k, v in example.items():
                    keywords = v['keywords'].split(', ')
                    for word in keywords:
                        if word in label2idx:
                            split_nonzero_classes.append(word)

        split_nonzero_classes = list(set(split_nonzero_classes))
        split_zero_classes = list(set(all_indices) - set(split_nonzero_classes))
        all_split_zero_classes.extend(split_zero_classes)
    
    all_split_zero_classes = list(set(all_split_zero_classes))
    # print(all_split_zero_classes)
    all_split_nonzero_classes = list(set(all_indices) - set(all_split_zero_classes))

    new_label2idx = {}
    new_idx2label = {}
    for i, nonzero in enumerate(all_split_nonzero_classes):
        keyword = idx2label[str(nonzero)]
        new_label2idx[keyword] = i
        new_idx2label[i] = keyword

    with open(f'filtered_label2idx.json', 'w') as f3:
        json.dump(new_label2idx, f3)

    with open(f'filtered_idx2label.json', 'w') as f4:
        json.dump(new_idx2label, f4)


def filter_split():
    """
    from filtered labels, filter splits
    """
    with open(f'intersection_label2idx.json', 'r') as f1:
        label2idx = json.load(f1)
    with open(f'intersection_idx2label.json', 'r') as f2:
        idx2label = json.load(f2)
    datapath_list = ['DeepEyeNet_train.json', 'DeepEyeNet_valid.json', 'DeepEyeNet_test.json']
    for i, file in enumerate(datapath_list):
        filtered_split_list = []
        with open(file, 'r') as f3:
            split_list = json.load(f3)
            for example in split_list:
                for k, v in example.items():
                    keywords = v['keywords']
                    if keywords in label2idx:
                        filtered_split_list.append(example)

        with open(f'filtered_{file}', 'w') as f4:
            json.dump(filtered_split_list, f4)


def filter_split_dict():
    """
    build 3 splits' count dict, filter < 8, then intersect
    """
    datapath_list = ['DeepEyeNet_train.json', 'DeepEyeNet_valid.json', 'DeepEyeNet_test.json']
    word2count_list = []  # 3 splits word2count dict
    for file in datapath_list:
        word2count = {}  # word: count
        with open(file, 'r') as f1:
            split_list = json.load(f1)
            for example in split_list:
                for k, v in example.items():
                    keywords = v['keywords'].split(', ')
                    for word in keywords:
                        if word in word2count:
                            word2count[word] = word2count[word] + 1
                        else:
                            word2count[word] = 1
        word_to_delete = []
        for word, count in word2count.items():
            if count < 8:
                word_to_delete.append(word)
        for less_word in word_to_delete:
            word2count.pop(less_word)
        word2count_list.append(word2count)
    
    words_set_list = [set(list(item.keys())) for item in word2count_list]
    inter = set.intersection(*words_set_list)
    least8_label2idx = {}
    least8_idx2label = {}
    for i, word in enumerate(inter):
        least8_label2idx[word] = i
        least8_idx2label[i] = word
    
    with open('least8_label2idx.json', 'w') as f2:
        json.dump(least8_label2idx, f2)
    with open('least8_idx2label.json', 'w') as f3:
        json.dump(least8_idx2label, f3)


def filter_splits_least8():
    """
    filter 3 splits by at least 8 label dicts
    """
    with open('least8_label2idx.json', 'r') as f1:
        label2idx = json.load(f1)
    # with open('least8_idx2label.json', 'r') as f2:
    #     idx2label = json.load(f2)
    
    datapath_list = ['DeepEyeNet_train.json', 'DeepEyeNet_valid.json', 'DeepEyeNet_test.json']

    for file in datapath_list:
        filtered_list = []
        with open(file, 'r') as f3:
            split_list = json.load(f3)
            for example in split_list:
                for k, v in example.items():
                    keywords = v['keywords'].split(', ')
                    if any(item in label2idx for item in keywords):
                        filtered_list.append(example)
        
        with open(f'filtered_{file}', 'w') as f2:
            json.dump(filtered_list, f2)


if __name__ == '__main__':
    # main()
    # build_label_dict()
    # filter_zero_classes()
    # filter_split()
    # filter_split_dict()
    filter_splits_least8()
